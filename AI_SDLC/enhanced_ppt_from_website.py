"""
Enhanced PowerPoint Presentation: Making NXOP AI-Native
Adds visual elements to better match the Reveal.js UI
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Brand Colors
AA_RED = RGBColor(200, 10, 40)
AA_DARK_BLUE = RGBColor(0, 75, 135)
AA_LIGHT_BLUE = RGBColor(0, 120, 210)
AA_SILVER = RGBColor(167, 170, 173)
AA_DARK_GRAY = RGBColor(43, 43, 43)
SUCCESS_GREEN = RGBColor(40, 167, 69)
WARNING_ORANGE = RGBColor(255, 140, 0)

CARD_BG = RGBColor(248, 249, 250)

# Helper functions
def add_icon(shape, icon, size=36, color=AA_DARK_BLUE):
    """Add an emoji icon to a shape (cross-platform)"""
    tf = shape.text_frame
    tf.text = icon
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

def add_card(slide, left, top, width, height, title, content, icon=None, color=AA_DARK_BLUE):
    """Add a card-style rectangle with optional icon and text"""
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = color
    card.line.width = Pt(2)
    # Shadows are not fully supported in python-pptx; skip color assignment
    card.shadow.inherit = False
    card.shadow.blur_radius = 8
    card.shadow.distance = 4
    # card.shadow.color.rgb = AA_SILVER  # Not supported
    # card.shadow.transparency = 0.5
    
    tf = card.text_frame
    tf.clear()
    if icon:
        p = tf.add_paragraph()
        p.text = icon
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = content
    p.font.size = Pt(14)
    p.font.color.rgb = AA_DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    return card

def add_progress_bar(slide, left, top, width, height, percent, label, color=SUCCESS_GREEN):
    """Add a static progress bar with label"""
    # Background bar
    bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = AA_SILVER
    bar_bg.line.fill.background()
    # Foreground bar
    bar_fg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, int(width * percent), height)
    bar_fg.fill.solid()
    bar_fg.fill.fore_color.rgb = color
    bar_fg.line.fill.background()
    # Label
    label_box = slide.shapes.add_textbox(left + width + 0.05*width, top, Inches(1.2), height)
    tf = label_box.text_frame
    tf.text = f"{int(percent*100)}% {label}"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = AA_DARK_BLUE
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "Making NXOP AI-Native"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.text = "A Path to Speed, Reliability, and Scale"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = AA_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    # Meta
    meta_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    tf = meta_box.text_frame
    tf.text = "NXOP Program | Technology Leadership | January 9, 2026"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = AA_SILVER
    p.alignment = PP_ALIGN.CENTER

def add_agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Agenda")
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    agenda_items = [
        ("1. Current Enterprise Metrics", "Where Time Goes in SDLC (5 mins)", "📊"),
        ("2. AI-Native SDLC Adoption", "Crawl, Walk, Run Model (5 mins)", "🤖"),
        ("3. Live Demo", "MCP + AI in Action (5 mins)", "🧑‍💻"),
        ("4. Progress Dashboard", "Executive Overview (3 mins)", "📈"),
        ("5. Business Outcomes", "Expected ROI & Impact (4 mins)", "🎯")
    ]
    for title, desc, icon in agenda_items:
        p = tf.add_paragraph()
        p.text = f"{icon} {title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = AA_DARK_BLUE
        p.space_before = Pt(10)
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        p2.font.color.rgb = AA_DARK_GRAY
        p2.level = 1

def add_metrics_cards_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Current Enterprise Metrics"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = AA_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    # Cards
    add_card(slide, Inches(0.5), Inches(1.5), Inches(3), Inches(2), "5.9 Days", "Per Feature", "📅", AA_LIGHT_BLUE)
    add_card(slide, Inches(3.7), Inches(1.5), Inches(3), Inches(2), "26 Mins", "To Detect Issues", "⏱", WARNING_ORANGE)
    add_card(slide, Inches(6.9), Inches(1.5), Inches(3), Inches(2), "52 Days", "Avg. Dwell Time", "⌛", SUCCESS_GREEN)
    # Manual Process
    manual_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.2))
    tf = manual_box.text_frame
    tf.text = "Manual Process Reality: Requirements, Development, Testing, Deployment all have bottlenecks."
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = AA_DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    # Critical Finding
    crit_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.8))
    tf = crit_box.text_frame
    tf.text = "⚠ NXOP cannot scale without transforming delivery."
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = AA_RED
    p.alignment = PP_ALIGN.CENTER

def add_progress_dashboard_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Executive Progress Dashboard"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = AA_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    # Progress Bar
    add_progress_bar(slide, Inches(1.5), Inches(1.5), Inches(7), Inches(0.5), 0.35, "Overall AI-Native SDLC Adoption", AA_LIGHT_BLUE)
    # Initiative Status
    add_card(slide, Inches(1.5), Inches(2.3), Inches(2.5), Inches(1.2), "Copilot", "Completed", "✅", SUCCESS_GREEN)
    add_card(slide, Inches(4.1), Inches(2.3), Inches(2.5), Inches(1.2), "MCP", "In Progress", "🔄", WARNING_ORANGE)
    add_card(slide, Inches(6.7), Inches(2.3), Inches(2.5), Inches(1.2), "AI Agents", "Planned", "🕒", AA_LIGHT_BLUE)
    # Metrics
    add_card(slide, Inches(1.5), Inches(3.7), Inches(2.5), Inches(1.2), "40%", "Velocity Improvement", "⚡", SUCCESS_GREEN)
    add_card(slide, Inches(4.1), Inches(3.7), Inches(2.5), Inches(1.2), "3/10", "Tools Connected via MCP", "🔌", WARNING_ORANGE)
    add_card(slide, Inches(6.7), Inches(3.7), Inches(2.5), Inches(1.2), "75%", "Test Coverage Increase", "🧪", AA_DARK_BLUE)

def set_slide_title(slide, title_text):
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = AA_RED

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    add_title_slide(prs)
    add_agenda_slide(prs)
    add_metrics_cards_slide(prs)
    add_progress_dashboard_slide(prs)

    # CRAWL Phase
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 255, 240)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "AI Maturity: CRAWL Phase"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER
    add_card(slide, Inches(1), Inches(1.5), Inches(2.7), Inches(2), "🛠 Tools Deployed", "GitHub Copilot\nChat-based AI", None, SUCCESS_GREEN)
    add_card(slide, Inches(3.8), Inches(1.5), Inches(2.7), Inches(2), "⚡ What It Does", "Code generation\nDocumentation help\nDebugging suggestions", None, SUCCESS_GREEN)
    add_card(slide, Inches(6.6), Inches(1.5), Inches(2.7), Inches(2), "📈 Value Gained", "Faster development\nQuick wins\nLower learning curve", None, SUCCESS_GREEN)
    crit_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    tf = crit_box.text_frame
    tf.text = "⚠ Key Limitation: Disconnected from NXOP systems, architecture standards, and vendor contracts."
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = AA_RED
    p.alignment = PP_ALIGN.CENTER

    # WALK Phase
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 250, 240)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "AI Maturity: WALK Phase"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WARNING_ORANGE
    p.alignment = PP_ALIGN.CENTER
    add_card(slide, Inches(1), Inches(1.5), Inches(2.7), Inches(2), "👥 Role-Based AI", "Developer assistant\nTest automation\nSRE triage", None, WARNING_ORANGE)
    add_card(slide, Inches(3.8), Inches(1.5), Inches(2.7), Inches(2), "🔗 Connected To", "Vendor specs\nCI/CD pipelines\nMetrics & logs", None, WARNING_ORANGE)
    add_card(slide, Inches(6.6), Inches(1.5), Inches(2.7), Inches(2), "🚀 Impact", "Weeks → Days\n35% → 75% test coverage\nFewer regressions", None, WARNING_ORANGE)
    enabler_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    tf = enabler_box.text_frame
    tf.text = "✨ Key Enabler: MCP connects vendor products for unified development intelligence."
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = AA_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # RUN Phase
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "AI Maturity: RUN Phase"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = AA_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    add_card(slide, Inches(1), Inches(1.5), Inches(2.7), Inches(2), "🤖 AI Agents In", "Delivery workflows\nSRE operations\nChange management", None, AA_DARK_BLUE)
    add_card(slide, Inches(3.8), Inches(1.5), Inches(2.7), Inches(2), "👥 Hybrid Teams", "Human + AI squads\nCollaborative intelligence\nContinuous learning loops", None, AA_DARK_BLUE)
    add_card(slide, Inches(6.6), Inches(1.5), Inches(2.7), Inches(2), "⭐ Outcomes", "Predictive reliability\nSelf-optimizing ops\nFull platform autonomy", None, AA_DARK_BLUE)
    trans_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    tf = trans_box.text_frame
    tf.text = "👑 Key Transformation: NXOP operates as a self-improving digital platform with autonomous agents."
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER

    # DEMO Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 249, 250)
    demo_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(2))
    tf = demo_box.text_frame
    tf.text = "DEMO"
    p = tf.paragraphs[0]
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = AA_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.text = "See MCP + AI in Action"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.color.rgb = AA_DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    features_box = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(6), Inches(1))
    tf = features_box.text_frame
    tf.text = "Vendor Integration • Code Generation • AI Testing"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = AA_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Business Outcomes
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 255)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Expected Business Outcomes"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = AA_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    add_card(slide, Inches(1), Inches(1.5), Inches(2.7), Inches(2), "🚀 Outcome-Driven Dev", "AI handles undifferentiated heavy lifting. Engineers focus on business value.", None, AA_LIGHT_BLUE)
    add_card(slide, Inches(3.8), Inches(1.5), Inches(2.7), Inches(2), "🤝 Multi-Vendor Integration", "MCP reduces onboarding from months to weeks.", None, AA_LIGHT_BLUE)
    add_card(slide, Inches(6.6), Inches(1.5), Inches(2.7), Inches(2), "🛡 Risk Mitigation", "Automated compliance checking and unified experience.", None, AA_LIGHT_BLUE)
    add_card(slide, Inches(1), Inches(3.8), Inches(2.7), Inches(1.2), "💡 Developer Excellence", "Attract and retain top talent.", None, AA_DARK_BLUE)
    add_card(slide, Inches(3.8), Inches(3.8), Inches(2.7), Inches(1.2), "⚡ Business Agility", "Faster time-to-market for new airline capabilities.", None, AA_DARK_BLUE)
    add_card(slide, Inches(6.6), Inches(3.8), Inches(2.7), Inches(1.2), "🏆 Organizational Maturity", "AI-native enterprise with connected systems.", None, AA_DARK_BLUE)

    # Impact Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 255, 255)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Expected Impact"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER
    add_card(slide, Inches(1.5), Inches(1.5), Inches(2.5), Inches(1.5), "1+ Year", "Timeline Reduction", "📅", SUCCESS_GREEN)
    add_card(slide, Inches(4.1), Inches(1.5), Inches(2.5), Inches(1.5), "40%", "Productivity Gain", "🏆", SUCCESS_GREEN)
    add_card(slide, Inches(6.7), Inches(1.5), Inches(2.5), Inches(1.5), "Significant", "Cost Savings", "💲", SUCCESS_GREEN)

    # Thank You Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = AA_DARK_BLUE
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.text = "Questions & Discussion"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.text = "Let's Transform NXOP Together"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = AA_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    return prs

if __name__ == "__main__":
    print("Generating enhanced PowerPoint presentation...")
    presentation = create_presentation()
    output_path = "presentations/ai-enabled-sdlc-nxop/NXOP_AI_Native_Enhanced_Presentation.pptx"
    presentation.save(output_path)
    print(f"✓ Enhanced presentation saved to: {output_path}")
