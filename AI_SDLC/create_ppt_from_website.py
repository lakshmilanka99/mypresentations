"""
Generate PowerPoint Presentation: Making NXOP AI-Native
Based on the Reveal.js presentation website
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# American Airlines Brand Colors
AA_RED = RGBColor(200, 10, 40)
AA_DARK_BLUE = RGBColor(0, 75, 135)
AA_LIGHT_BLUE = RGBColor(0, 120, 210)
AA_SILVER = RGBColor(167, 170, 173)
AA_DARK_GRAY = RGBColor(43, 43, 43)
SUCCESS_GREEN = RGBColor(40, 167, 69)
WARNING_ORANGE = RGBColor(255, 140, 0)

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = AA_DARK_BLUE
    
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "Making NXOP AI-Native"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.text = "A Path to Speed, Reliability, and Scale"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = AA_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Meta info
    meta_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    tf = meta_box.text_frame
    tf.text = "NXOP Program | Technology Leadership | January 9, 2026"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = AA_SILVER
    p.alignment = PP_ALIGN.CENTER

def add_agenda_slide(prs):
    """Slide 2: Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Agenda")
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    agenda_items = [
        ("1. Current Enterprise Metrics", "Where Time Goes in SDLC (5 mins)"),
        ("2. AI-Native SDLC Adoption", "Crawl, Walk, Run Model (5 mins)"),
        ("3. Live Demo", "MCP + AI in Action (5 mins)"),
        ("4. Progress Dashboard", "Executive Overview (3 mins)"),
        ("5. Business Outcomes", "Expected ROI & Impact (4 mins)")
    ]
    
    for title, desc in agenda_items:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = AA_DARK_BLUE
        p.space_before = Pt(10)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        p2.font.color.rgb = AA_DARK_GRAY
        p2.level = 1

def add_current_state_slide(prs):
    """Slide 3: Current Enterprise Metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Current Enterprise Metrics")
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Key Metrics
    metrics = [
        "5.9 Days per Feature",
        "26 Minutes to Detect Issues",
        "52 Days Average Dwell Time"
    ]
    
    p = tf.add_paragraph()
    p.text = "Key Performance Indicators:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = AA_DARK_BLUE
    
    for metric in metrics:
        p = tf.add_paragraph()
        p.text = metric
        p.font.size = Pt(18)
        p.font.color.rgb = SUCCESS_GREEN
        p.font.bold = True
        p.level = 1
    
    # Manual Process Reality
    p = tf.add_paragraph()
    p.text = "\nManual Process Reality Across SDLC:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = AA_DARK_BLUE
    p.space_before = Pt(20)
    
    processes = [
        "Requirements: Manual code analysis, reverse-engineering, tribal knowledge",
        "Development: Line-by-line translation, manual boilerplate",
        "Testing: Manual test creation and verification",
        "Deployment: Manual planning and reactive incident response"
    ]
    
    for process in processes:
        p = tf.add_paragraph()
        p.text = process
        p.font.size = Pt(16)
        p.level = 1
    
    # Critical Finding
    p = tf.add_paragraph()
    p.text = "\n⚠ Critical Finding:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = AA_RED
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "NXOP cannot scale without transforming delivery"
    p.font.size = Pt(18)
    p.font.color.rgb = AA_DARK_GRAY
    p.level = 1

def add_sdlc_time_slide(prs):
    """Slide 4: Where Time Goes in SDLC"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Where Time Goes in SDLC")
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "SDLC Stages:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = AA_DARK_BLUE
    
    stages = [
        "Requirements • Design • Development",
        "Testing • Deployment • Monitoring",
        "Maintenance • Planning"
    ]
    
    for stage in stages:
        p = tf.add_paragraph()
        p.text = stage
        p.font.size = Pt(18)
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n⏱ Most time spent waiting between stages:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WARNING_ORANGE
    p.space_before = Pt(20)
    
    wait_times = [
        "Environment setup",
        "Code reviews",
        "Testing",
        "Approvals",
        "Handoffs"
    ]
    
    for item in wait_times:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 1
    
    # AI Solution
    p = tf.add_paragraph()
    p.text = "\n🤖 AI Solution:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = AA_LIGHT_BLUE
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Eliminate meetings • Automate handoffs • Real-time coordination"
    p.font.size = Pt(18)
    p.level = 1

def add_maturity_crawl_slide(prs):
    """Slide 5: CRAWL Phase"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "AI Maturity: CRAWL Phase")
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("🛠 Tools Deployed", ["GitHub Copilot", "Chat-based AI"]),
        ("⚡ What It Does", ["Code generation", "Documentation help", "Debugging suggestions"]),
        ("📈 Value Gained", ["Faster development", "Quick wins", "Lower learning curve"])
    ]
    
    for title, items in sections:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = SUCCESS_GREEN
        p.space_before = Pt(15)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.level = 1
    
    # Limitation
    p = tf.add_paragraph()
    p.text = "\n⚠ Key Limitation:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = AA_RED
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Disconnected from NXOP systems, architecture standards, and vendor contracts"
    p.font.size = Pt(16)
    p.level = 1

def add_maturity_walk_slide(prs):
    """Slide 6: WALK Phase"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "AI Maturity: WALK Phase")
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("👥 Role-Based AI", ["Developer assistant", "Test automation", "SRE triage"]),
        ("🔗 Connected To", ["Vendor specs", "CI/CD pipelines", "Metrics & logs"]),
        ("🚀 Impact", ["Weeks → Days", "35% → 75% test coverage", "Fewer regressions"])
    ]
    
    for title, items in sections:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WARNING_ORANGE
        p.space_before = Pt(15)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.level = 1
    
    # Key Enabler
    p = tf.add_paragraph()
    p.text = "\n✨ Key Enabler: MCP"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = AA_LIGHT_BLUE
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Vendor products connected through Model Context Protocol for unified development intelligence"
    p.font.size = Pt(16)
    p.level = 1

def add_maturity_run_slide(prs):
    """Slide 7: RUN Phase"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "AI Maturity: RUN Phase")
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("🤖 AI Agents In", ["Delivery workflows", "SRE operations", "Change management"]),
        ("👥 Hybrid Teams", ["Human + AI squads", "Collaborative intelligence", "Continuous learning loops"]),
        ("⭐ Outcomes", ["Predictive reliability", "Self-optimizing ops", "Full platform autonomy"])
    ]
    
    for title, items in sections:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = AA_DARK_BLUE
        p.space_before = Pt(15)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.level = 1
    
    # Key Transformation
    p = tf.add_paragraph()
    p.text = "\n👑 Key Transformation:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "NXOP operates as a self-improving digital platform with autonomous agents"
    p.font.size = Pt(16)
    p.level = 1

def add_demo_slide(prs):
    """Slide 8: Live Demo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 249, 250)
    
    # DEMO text
    demo_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(2))
    tf = demo_box.text_frame
    tf.text = "DEMO"
    p = tf.paragraphs[0]
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = AA_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.text = "See MCP + AI in Action"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.color.rgb = AA_DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Features
    features_box = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(6), Inches(1))
    tf = features_box.text_frame
    tf.text = "Vendor Integration • Code Generation • AI Testing"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = AA_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

def add_progress_dashboard_slide(prs):
    """Slide 9: Executive Progress Dashboard"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Executive Progress Dashboard")
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Current Status
    p = tf.add_paragraph()
    p.text = "ℹ Current Status:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = AA_LIGHT_BLUE
    
    p = tf.add_paragraph()
    p.text = "NXOP is transitioning from CRAWL → WALK phase"
    p.font.size = Pt(18)
    p.level = 1
    
    # Overall Progress
    p = tf.add_paragraph()
    p.text = "\n📊 Overall AI-Native SDLC Adoption: 35%"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = AA_DARK_BLUE
    p.space_before = Pt(20)
    
    # Initiative Status
    p = tf.add_paragraph()
    p.text = "\n✅ Completed:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.space_before = Pt(15)
    
    completed = ["GitHub Copilot Deployment", "Team Training"]
    for item in completed:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n🔄 In Progress:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WARNING_ORANGE
    p.space_before = Pt(10)
    
    in_progress = ["MCP Integration", "Test Automation", "CI/CD Enhancement"]
    for item in in_progress:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n📅 Planned:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = AA_LIGHT_BLUE
    p.space_before = Pt(10)
    
    planned = ["AI Agents", "Predictive Operations"]
    for item in planned:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 1

def add_metrics_slide(prs):
    """Slide 10: Key Metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Key Performance Metrics")
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    metrics = [
        ("40%", "Velocity Improvement"),
        ("3/10", "Tools Connected via MCP"),
        ("75%", "Test Coverage Increase")
    ]
    
    for value, label in metrics:
        p = tf.add_paragraph()
        p.text = f"{value} - {label}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = AA_DARK_BLUE
        p.space_before = Pt(20)

def add_business_outcomes_slide(prs):
    """Slide 11: Expected Business Outcomes"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Expected Business Outcomes")
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Hero Statement
    p = tf.add_paragraph()
    p.text = "🚀 From Non-Differentiating Work to Outcome-Driven Development"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = AA_LIGHT_BLUE
    
    p = tf.add_paragraph()
    p.text = "Free developers from plumbing tasks — AI handles undifferentiated heavy lifting"
    p.font.size = Pt(18)
    p.level = 1
    p.space_after = Pt(20)
    
    # Strategic Benefits
    p = tf.add_paragraph()
    p.text = "Strategic Benefits:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = AA_DARK_BLUE
    p.space_before = Pt(15)
    
    benefits = [
        "Accelerated Multi-Vendor Integration: Months → Weeks",
        "Unified Development Experience: Single IDE with all vendor context",
        "Organizational Maturity: AI-native enterprise with connected systems",
        "Risk Mitigation: Automated compliance checking",
        "Developer Excellence: Attract and retain top talent",
        "Business Agility: Faster time-to-market for airline capabilities"
    ]
    
    for benefit in benefits:
        p = tf.add_paragraph()
        p.text = benefit
        p.font.size = Pt(16)
        p.level = 1

def add_final_metrics_slide(prs):
    """Slide 12: Final Impact Metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Expected Impact")
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    impacts = [
        ("1+ Year", "Timeline Reduction"),
        ("40%", "Productivity Gain"),
        ("Significant", "Cost Savings")
    ]
    
    for value, label in impacts:
        p = tf.add_paragraph()
        p.text = f"{label}:"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = AA_DARK_BLUE
        p.space_before = Pt(20)
        
        p = tf.add_paragraph()
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = SUCCESS_GREEN
        p.level = 1

def add_thank_you_slide(prs):
    """Slide 13: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = AA_DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.text = "Questions & Discussion"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.text = "Let's Transform NXOP Together"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = AA_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

def set_slide_title(slide, title_text):
    """Set the title for a slide"""
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = AA_RED

def create_presentation():
    """Create the complete PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add all slides
    add_title_slide(prs)
    add_agenda_slide(prs)
    add_current_state_slide(prs)
    add_sdlc_time_slide(prs)
    add_maturity_crawl_slide(prs)
    add_maturity_walk_slide(prs)
    add_maturity_run_slide(prs)
    add_demo_slide(prs)
    add_progress_dashboard_slide(prs)
    add_metrics_slide(prs)
    add_business_outcomes_slide(prs)
    add_final_metrics_slide(prs)
    add_thank_you_slide(prs)
    
    return prs

if __name__ == "__main__":
    print("Generating PowerPoint presentation from website content...")
    presentation = create_presentation()
    output_path = "presentations/ai-enabled-sdlc-nxop/NXOP_AI_Native_Presentation.pptx"
    presentation.save(output_path)
    print(f"✓ Presentation saved to: {output_path}")
    print(f"✓ Total slides: {len(presentation.slides)}")
