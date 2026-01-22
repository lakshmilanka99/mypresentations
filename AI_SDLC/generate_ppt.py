"""
Generate PowerPoint Presentation: Making NXOP AI-Native
Based on the Reveal.js presentation content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# American Airlines Brand Colors
AA_RED = RGBColor(200, 10, 40)
AA_DARK_BLUE = RGBColor(0, 75, 135)
AA_LIGHT_BLUE = RGBColor(0, 120, 210)
AA_SILVER = RGBColor(167, 170, 173)
AA_DARK_GRAY = RGBColor(43, 43, 43)
SUCCESS_GREEN = RGBColor(40, 167, 69)
WARNING_ORANGE = RGBColor(255, 140, 0)

def create_presentation():
    """Create the complete PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide - Hero
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_background(slide1, RGBColor(0, 51, 102), RGBColor(0, 102, 204))
    
    title = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title.text_frame
    title_frame.text = "Accelerating NXOP Delivery with AI Across the SDLC"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle = slide1.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(0.8))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "From requirements to reliability — improving speed, quality, and operational confidence"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = RGBColor(220, 220, 220)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem - Growing Demand & Complexity
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide2, "The Challenge: Growing Demand & Complexity")
    
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = "Current Pain Points:"
    
    points = [
        "Rapid integration growth across NXOP ecosystem",
        "Long feedback loops in testing and validation",
        "Manual testing bottlenecks slowing delivery",
        "Reactive operations leading to increased MTTR",
        "Vendor onboarding delays and inconsistencies"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(20)
    
    # Slide 3: Reframe - What AI in SDLC Really Means
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide3, "AI Across the SDLC: System-Level Productivity")
    
    content = slide3.placeholders[1]
    tf = content.text_frame
    tf.text = "AI Applications Throughout Development Lifecycle:"
    
    ai_areas = [
        ("Planning", "Requirements analysis, story generation, effort estimation"),
        ("Development", "Code generation, refactoring, documentation"),
        ("Testing", "Test case generation, automated regression, contract validation"),
        ("Release", "Risk scoring, deployment automation, rollback decisions"),
        ("Operations", "Incident triage, root cause analysis, runbook generation"),
        ("Maintenance", "Technical debt identification, dependency updates")
    ]
    
    for area, description in ai_areas:
        p = tf.add_paragraph()
        p.text = f"{area}: {description}"
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 4: Target Operating Model
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide4, "AI-Enabled Pipeline: Target Operating Model")
    
    content = slide4.placeholders[1]
    tf = content.text_frame
    tf.text = "Delivery Flow with AI Assistance:"
    
    stages = [
        "Requirements → AI-assisted story refinement",
        "Design → Automated architecture validation",
        "Development → Copilot-accelerated coding",
        "Testing → AI-generated test suites",
        "Deployment → Risk-scored releases",
        "Operations → Intelligent incident response"
    ]
    
    for stage in stages:
        p = tf.add_paragraph()
        p.text = stage
        p.level = 1
        p.font.size = Pt(20)
    
    # Slide 5: NXOP Priority Use Cases
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide5, "Priority Use Cases for NXOP Platform")
    
    content = slide5.placeholders[1]
    tf = content.text_frame
    
    use_cases = [
        ("Integration Development Acceleration", 
         "Faster API integration development with code generation and documentation"),
        ("Automated Contract Testing",
         "AI-generated contract tests ensuring integration compatibility"),
        ("Regression Test Automation",
         "Comprehensive test coverage with minimal manual effort"),
        ("Incident Triage & Observability",
         "Intelligent log analysis and anomaly detection"),
        ("Runbook & Remediation Automation",
         "Auto-generated response procedures and proactive fixes")
    ]
    
    for title, desc in use_cases:
        p = tf.add_paragraph()
        p.text = f"{title}"
        p.font.bold = True
        p.font.size = Pt(18)
        p.level = 0
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        p2.level = 1
    
    # Slide 6: Impact & Metrics
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide6, "Expected Impact & Key Metrics")
    
    content = slide6.placeholders[1]
    tf = content.text_frame
    tf.text = "Measurable Improvements:"
    
    metrics = [
        "Cycle Time: 40-50% reduction in feature delivery time",
        "Deployment Frequency: 2-3x increase in safe deployments",
        "Test Coverage: 80%+ automated vs 40% manual baseline",
        "MTTR: 30-40% faster incident resolution",
        "Vendor Onboarding: 50% reduction in integration time"
    ]
    
    for metric in metrics:
        p = tf.add_paragraph()
        p.text = metric
        p.level = 1
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 102, 51)
    
    # Slide 7: Governance & Risk Controls
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide7, "Governance & Responsible AI Controls")
    
    content = slide7.placeholders[1]
    tf = content.text_frame
    tf.text = "Essential Guardrails:"
    
    controls = [
        "Human Approval: Critical decisions require human validation",
        "Audit Trails: Complete logging of AI-assisted changes",
        "Security Scanning: Automated vulnerability detection",
        "Policy Compliance: Regulatory and corporate policy enforcement",
        "Code Review: Mandatory peer review for AI-generated code",
        "Data Privacy: Restricted access to sensitive information"
    ]
    
    for control in controls:
        p = tf.add_paragraph()
        p.text = control
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 8: Vendor & Partner Alignment
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide8, "Vendor & Partner Alignment")
    
    content = slide8.placeholders[1]
    tf = content.text_frame
    tf.text = "AI-Enabled Co-Build Expectations:"
    
    expectations = [
        "Vendors must provide comprehensive API documentation",
        "Standardized telemetry and observability integration",
        "Test automation as part of delivery acceptance",
        "AI-compatible code quality and documentation standards",
        "Benefits: Faster onboarding, consistency, lower integration risk"
    ]
    
    for expectation in expectations:
        p = tf.add_paragraph()
        p.text = expectation
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 9: Implementation Roadmap
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide9, "Implementation Roadmap")
    
    content = slide9.placeholders[1]
    tf = content.text_frame
    
    phases = [
        ("Phase 1: Foundation (Q1-Q2 2026)",
         "GitHub Copilot rollout, AI-assisted test generation, developer training"),
        ("Phase 2: Pipeline Intelligence (Q3-Q4 2026)",
         "Automated risk scoring, CI/CD optimization, observability AI"),
        ("Phase 3: Predictive Operations (2027)",
         "Proactive incident prevention, self-healing systems, predictive capacity planning")
    ]
    
    for phase, details in phases:
        p = tf.add_paragraph()
        p.text = phase
        p.font.bold = True
        p.font.size = Pt(20)
        p.level = 0
        
        p2 = tf.add_paragraph()
        p2.text = details
        p2.font.size = Pt(16)
        p2.level = 1
    
    # Slide 10: Strategic Value to NXOP
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide10, "Strategic Value to NXOP Platform")
    
    content = slide10.placeholders[1]
    tf = content.text_frame
    tf.text = "Business-Aligned Outcomes:"
    
    values = [
        "Scalability: Handle growing integration demands efficiently",
        "Reliability: Higher quality through comprehensive testing",
        "Vendor Independence: Reduce lock-in through automation",
        "Regulatory Response: Faster adaptation to compliance changes",
        "Operational Excellence: Shift from reactive to proactive",
        "Cost Efficiency: Optimize resource utilization and reduce rework"
    ]
    
    for value in values:
        p = tf.add_paragraph()
        p.text = value
        p.level = 1
        p.font.size = Pt(20)
    
    # Slide 11: Call to Action
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_background(slide11, RGBColor(0, 102, 51), RGBColor(0, 153, 76))
    
    cta_title = slide11.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    cta_title_frame = cta_title.text_frame
    cta_title_frame.text = "Next Steps"
    cta_title_para = cta_title_frame.paragraphs[0]
    cta_title_para.font.size = Pt(48)
    cta_title_para.font.bold = True
    cta_title_para.font.color.rgb = RGBColor(255, 255, 255)
    cta_title_para.alignment = PP_ALIGN.CENTER
    
    cta_content = slide11.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(3))
    cta_tf = cta_content.text_frame
    
    actions = [
        "Align on AI-enabled SDLC as our operating model",
        "Include AI expectations in vendor contracts",
        "Fund automation alongside feature delivery",
        "Commit to measurement and continuous improvement"
    ]
    
    for action in actions:
        p = cta_tf.add_paragraph()
        p.text = f"✓ {action}"
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(12)
    
    return prs

def set_slide_title(slide, title_text):
    """Set the title for a slide"""
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def add_gradient_background(slide, color1, color2):
    """Add a gradient background to a slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

if __name__ == "__main__":
    print("Generating PowerPoint presentation...")
    presentation = create_presentation()
    output_path = "presentations/ai-enabled-sdlc-nxop/AI_SDLC_NXOP_Presentation.pptx"
    presentation.save(output_path)
    print(f"✓ Presentation saved to: {output_path}")
